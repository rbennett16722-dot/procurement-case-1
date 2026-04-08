"""
Build v2.0 recommendations doc and slideshow for MSBX 5470 Case 1.
"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')

# ─── WORD DOCUMENT ───────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

doc = Document()

def bold_row(row):
    for cell in row.cells:
        for para in cell.paragraphs:
            for run in para.runs:
                run.bold = True

# ── Title block
doc.add_heading("MSBX 5470 — Case 1: Packaging & Dunnage", 0)
sub = doc.add_paragraph("Data Quality Review & Procurement Recommendations  |  Version 2.0")
sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
sub.runs[0].italic = True
doc.add_paragraph("Group 2  |  Specialty Aluminum Can Manufacturing  |  April 2026")
doc.add_paragraph("")

# ── Section 1: Executive Summary
doc.add_heading("1. Executive Summary", 1)
doc.add_paragraph(
    "This document presents updated findings from the Case 1 Packaging & Dunnage spend analysis. "
    "Version 2.0 corrects three confirmed misclassifications from the original 33-row PACKAGING dataset, "
    "identifies five additional suspected misclassifications, resolves six blank sub-category (Cat3) gaps, "
    "and provides strategic procurement recommendations based on the cleaned data."
)
doc.add_paragraph("")

doc.add_heading("Key Statistics (v2.0 Corrected Dataset)", 2)
t = doc.add_table(rows=1, cols=2)
t.style = "Table Grid"
bold_row(t.rows[0])
t.rows[0].cells[0].text = "Metric"
t.rows[0].cells[1].text = "Value"
for k, v in [
    ("Total Packaging Transactions (corrected)", "30  (was 33)"),
    ("Total Packaging Spend (corrected)", "$6,299,708.55  (was $6,373,081.88)"),
    ("Spend Removed via Confirmed Reclassification", "$73,373.33"),
    ("Additional Spend Under Review (5 suspects)", "~$429,071"),
    ("Unique Suppliers (corrected)", "22  (was 25)"),
    ("Avg Transaction Size", "$209,990"),
    ("% of Total Dataset Spend", "4.8%"),
    ("Avg Payment Terms", "~70 days"),
    ("Largest Sub-Category", "Layer Pads  (30.3%,  $1.91M)"),
    ("Top Supplier", "Angleboard  (23.2%,  $1.46M)"),
    ("EMEA Share of Spend", "52.1%"),
]:
    r = t.add_row()
    r.cells[0].text = k
    r.cells[1].text = v
doc.add_paragraph("")

# ── Section 2: Confirmed reclassifications
doc.add_heading("2. Confirmed Reclassifications (Applied in Workbook)", 1)
doc.add_paragraph(
    "Three transactions originally coded as PACKAGING were confirmed as misclassifications after PO description "
    "review. Corrections are recorded in the workbook (Spend File, Column L — Corrected Cat 2)."
)
t2 = doc.add_table(rows=1, cols=6)
t2.style = "Table Grid"
bold_row(t2.rows[0])
for i, h in enumerate(["Row", "Supplier", "PO Description", "Original", "Corrected To", "Spend"]):
    t2.rows[0].cells[i].text = h
for row_data in [
    ("543", "SODEXO", "2021 NOV - WORKWEAR", "PACKAGING", "HUMAN RESOURCES (HR)", "$59,402.56"),
    ("755", "JEBLA SA DE CV", "GAS MONTACARGAS 1-15 OCT", "PACKAGING", "ENERGY & UTILITIES", "$10,832.31"),
    ("211", "ARDAGH METAL BEVERAGE", "SHAFT, 265018", "PACKAGING", "EQUIPMENT - GENERAL", "$3,138.46"),
]:
    r = t2.add_row()
    for i, v in enumerate(row_data):
        r.cells[i].text = v
doc.add_paragraph("")

# ── Section 3: Additional suspects
doc.add_heading("3. Additional Suspected Misclassifications — Recommended for Review", 1)
doc.add_paragraph(
    "Five additional rows remain in the corrected PACKAGING dataset but appear misclassified based on PO "
    "description analysis. These were NOT automatically reclassified — team or instructor confirmation is needed. "
    "Combined potential impact: ~$429,071 (6.8% of corrected packaging spend)."
)
t3 = doc.add_table(rows=1, cols=6)
t3.style = "Table Grid"
bold_row(t3.rows[0])
for i, h in enumerate(["Row", "Supplier", "PO Description", "Current Cat 3", "Suspected Category", "Spend"]):
    t3.rows[0].cells[i].text = h
for row_data in [
    ("653", "EXTENDATA SOLUTIONS", "SSO FOR BUSINESS USERS", "PACKAGING LABELS", "IT / PROF. SERVICES", "$340,937.67"),
    ("643", "EXTENDATA SOLUTIONS", "UIM SERVICE RESPONSE", "PACKAGING LABELS", "IT / PROF. SERVICES", "$57,834.64"),
    ("579", "FOKSAN DOO", "FOTOKOPIR PAPIR A4 ALLROUND", "PAPER SHEET", "OFFICE SUPPLIES", "$16,650.98"),
    ("551", "NPB AUTOMATION", "CIRCUIT BREAKER, THERM-MAGN;5A", "WOODEN FRAME", "EQUIPMENT - GENERAL", "$9,938.46"),
    ("84",  "GRUPMICROS SMART SOLUTIONS", "RENO CONTRACTE MANT. MC92XX 1Y", "PACKAGING LABELS", "IT / PROF. SERVICES", "$3,710.32"),
]:
    r = t3.add_row()
    for i, v in enumerate(row_data):
        r.cells[i].text = v
doc.add_paragraph("")
p = doc.add_paragraph("")
p.add_run("Why these are suspect:").bold = True
doc.add_paragraph(
    "- EXTENDATA SOLUTIONS (Rows 643 & 653): Extendata is a software/IT company. "
    "'SSO FOR BUSINESS USERS' is Single Sign-On — a cybersecurity/IT service. "
    "'UIM SERVICE RESPONSE' is a service desk ticket response. Neither is a packaging label. "
    "Combined $398,772 classified as PACKAGING LABELS makes this the highest-impact flag."
)
doc.add_paragraph(
    "- FOKSAN DOO (Row 579): 'Fotokopir papir A4 Allround Bisniss' is Serbian/Croatian for 'photocopy paper A4.' "
    "This is office stationery, not packaging paper stock."
)
doc.add_paragraph(
    "- NPB AUTOMATION (Row 551): 'Circuit Breaker, Therm-Magn;5A' is an electrical/mechanical component. "
    "It is coded to Category 3 = WOODEN FRAME, which is implausible. Likely belongs in EQUIPMENT - GENERAL."
)
doc.add_paragraph(
    "- GRUPMICROS SMART SOLUTIONS (Row 84): 'RENO CONTRACTE MANT. MC92XX 1Y' translates to a 1-year "
    "renewal of a maintenance contract for the MC92XX — a Motorola/Zebra handheld mobile computer. "
    "This is an IT asset maintenance cost, not a packaging label."
)
doc.add_paragraph("")

# ── Section 4: Blank Cat3
doc.add_heading("4. Remaining Category 3 Gaps (Blank Sub-Categories)", 1)
doc.add_paragraph(
    "Six transactions remain in PACKAGING with no Category 3 sub-category assigned, totaling $97,463 "
    "(1.5% of corrected spend). These were verified as legitimate packaging transactions but lack a sub-category. "
    "Recommended assignments based on PO description analysis:"
)
t4 = doc.add_table(rows=1, cols=4)
t4.style = "Table Grid"
bold_row(t4.rows[0])
for i, h in enumerate(["Supplier", "PO Description", "Spend", "Recommended Cat 3"]):
    t4.rows[0].cells[i].text = h
for row_data in [
    ("KARTONFABRIK POSTENDORF", "PEAK SEASON SURCHARGE", "$47,037.77", "Layer Pads / End Bags (surcharge on packaging supply)"),
    ("MULTI SERVICES GSTJ", "NAC-1007", "$16,657.17", "Unknown — supplier confirms packaging material"),
    ("ULINE", "H-4184GR TOLVA GRIS", "$12,565.78", "Stretch Wrap (hopper/bin for packaging line)"),
    ("MULTI SERVICES GSTJ", "BAL4840BFHT-01:", "$9,382.83", "Stretch Wrap (BAL prefix = bale/stretch wrap SKU)"),
    ("ULINE", "H-3687GR J JAUREGUI", "$6,843.81", "Stretch Wrap / Packaging Supplies (ULINE product)"),
    ("MULTI SERVICES GSTJ", "ETIQ. 4X2 TT VERT", "$4,975.57", "Packaging Labels (etiqueta = label in Spanish)"),
]:
    r = t4.add_row()
    for i, v in enumerate(row_data):
        r.cells[i].text = v
doc.add_paragraph("")

# ── Section 5: Spend summary
doc.add_heading("5. Corrected Spend Analysis", 1)

doc.add_heading("Sub-Category Breakdown", 2)
t5 = doc.add_table(rows=1, cols=4)
t5.style = "Table Grid"
bold_row(t5.rows[0])
for i, h in enumerate(["Sub-Category", "Transactions", "Total Spend", "% of Total"]):
    t5.rows[0].cells[i].text = h
for row_data in [
    ("LAYER PADS",         "3",  "$1,909,249", "30.3%"),
    ("END BAGS",           "2",  "$1,071,522", "17.0%"),
    ("WOOD PALLET",        "5",  "$920,994",   "14.6%"),
    ("STRETCH WRAP",       "1",  "$866,966",   "13.8%"),
    ("TOP FRAMES",         "3",  "$539,107",   "8.6%"),
    ("PACKAGING LABELS",   "5",  "$430,913",   "6.8%"),
    ("WOODEN FRAME",       "2",  "$408,159",   "6.5%"),
    ("(Blank Cat3)",       "6",  "$97,463",    "1.5%"),
    ("STRAP TAPE",         "1",  "$31,711",    "0.5%"),
    ("PAPER SHEET",        "1",  "$16,651",    "0.3%"),
    ("PLASTIC LAYER SHEET","1",  "$6,973",     "0.1%"),
    ("TOTAL",              "30", "$6,299,709", "100%"),
]:
    r = t5.add_row()
    for i, v in enumerate(row_data):
        r.cells[i].text = v
doc.add_paragraph("")

doc.add_heading("Business Unit Breakdown", 2)
t5b = doc.add_table(rows=1, cols=4)
t5b.style = "Table Grid"
bold_row(t5b.rows[0])
for i, h in enumerate(["Business Unit", "Transactions", "Total Spend", "% of Total"]):
    t5b.rows[0].cells[i].text = h
for row_data in [
    ("EMEA",      "13", "$3,283,703", "52.1%"),
    ("SA",        "3",  "$1,444,528", "22.9%"),
    ("NCA",       "9",  "$1,141,689", "18.1%"),
    ("Corporate", "2",  "$398,772",   "6.3%"),
    ("Aerosol",   "3",  "$31,016",    "0.5%"),
    ("TOTAL",     "30", "$6,299,709", "100%"),
]:
    r = t5b.add_row()
    for i, v in enumerate(row_data):
        r.cells[i].text = v
doc.add_paragraph("")

doc.add_heading("Top 10 Suppliers by Spend", 2)
t6 = doc.add_table(rows=1, cols=5)
t6.style = "Table Grid"
bold_row(t6.rows[0])
for i, h in enumerate(["Rank", "Supplier", "Transactions", "Total Spend", "Cumulative %"]):
    t6.rows[0].cells[i].text = h
for row_data in [
    ("1",  "ANGLEBOARD",              "2", "$1,459,249", "23.2%"),
    ("2",  "GSD VERPACKUNGEN",        "1", "$1,011,306", "39.2%"),
    ("3",  "GROUP O",                 "1", "$866,966",   "52.9%"),
    ("4",  "EGYPTIAN PALLETS",        "1", "$657,912",   "63.4%"),
    ("5",  "ORBIS",                   "2", "$650,000",   "73.7%"),
    ("6",  "EXTENDATA SOLUTIONS *",   "3", "$423,277",   "80.5%"),
    ("7",  "FINE ARTS DE MEXICO",     "1", "$398,221",   "86.8%"),
    ("8",  "REHRIG",                  "2", "$339,107",   "92.2%"),
    ("9",  "ALDHANA WOOD INDUSTRIES", "1", "$118,791",   "94.0%"),
    ("10", "JAY WOOD INDUSTRY",       "1", "$79,535",    "95.3%"),
]:
    r = t6.add_row()
    for i, v in enumerate(row_data):
        r.cells[i].text = v
doc.add_paragraph("* EXTENDATA SOLUTIONS classification is under review — see Section 3.")
doc.add_paragraph("")

doc.add_heading("Payment Terms Distribution", 2)
t7 = doc.add_table(rows=1, cols=4)
t7.style = "Table Grid"
bold_row(t7.rows[0])
for i, h in enumerate(["Payment Terms", "Transactions", "Total Spend", "% of Pkg"]):
    t7.rows[0].cells[i].text = h
for row_data in [
    ("15-30 days",  "2",  "$89,474",   "1.4%"),
    ("31-60 days",  "13", "$2,125,441", "33.7%"),
    ("61-90 days",  "12", "$3,708,501", "58.9%"),
    ("91-120 days", "1",  "$47,038",   "0.7%"),
    ("121+ days",   "2",  "$329,255",  "5.2%"),
    ("TOTAL",       "30", "$6,299,709", "100%"),
]:
    r = t7.add_row()
    for i, v in enumerate(row_data):
        r.cells[i].text = v
doc.add_paragraph(
    "Note: ANGLEBOARD and REHRIG have 125-day payment terms ($329,255 combined). "
    "While favorable for working capital, terms >90 days may signal supplier relationship risk."
)
doc.add_paragraph("")

# ── Section 6: Recommendations
doc.add_heading("6. Procurement Recommendations", 1)
recs = [
    ("CRITICAL", "Audit EXTENDATA SOLUTIONS ($423K Potential Misclassification)",
     "Three transactions totaling $423,277 are classified as PACKAGING LABELS but PO descriptions indicate "
     "IT/software services (SSO login system, UIM service response, hardware bracket kit). If confirmed as IT "
     "spend, the packaging budget is overstated by 6.7% and the Pareto ranking shifts materially. "
     "Escalate to Finance/IT immediately before finalizing the category study."),
    ("HIGH", "Investigate 4 Additional Misclassification Flags",
     "NPB AUTOMATION (circuit breaker, $9,938), FOKSAN DOO (photocopy paper, $16,651), and GRUPMICROS "
     "SMART SOLUTIONS (scanner maintenance contract, $3,710) appear incorrectly coded to PACKAGING. "
     "Total at risk: ~$30,300. Review PO documentation for each."),
    ("HIGH", "Assign Cat3 Sub-Categories to 6 Blank Rows ($97,463)",
     "Six packaging transactions lack Category 3 sub-category, making them invisible in sub-category reporting. "
     "Recommended assignments are provided in Section 4. KARTONFABRIK POSTENDORF's $47,038 peak season "
     "surcharge is the largest single gap — this should be assigned to the primary sub-category of that supplier's "
     "other transactions (likely LAYER PADS or END BAGS)."),
    ("HIGH", "Address Supplier Concentration Risk",
     "Top 5 suppliers account for 73.7% of corrected packaging spend. Angleboard is the single source for "
     "layer pads ($1.46M, 23.2% of category). GSD Verpackungen is the only end bags supplier ($1.01M). "
     "A supply disruption at either would be critical. Run competitive bids for both sub-categories and "
     "develop a secondary-source qualification plan."),
    ("MEDIUM", "Consolidate Wood Pallet Supply Base",
     "Five separate suppliers provide WOOD PALLETS across $920,994. Consolidating to 2 preferred vendors "
     "would improve leverage and enable volume pricing. Recommended consolidation: retain Egyptian Pallets "
     "($657,912) as primary, qualify one secondary from Aldhana Wood ($118,791) or Jay Wood ($79,535). "
     "Evaluate ISPM-15 heat treatment capabilities for cross-border shipments."),
    ("MEDIUM", "Standardize Payment Terms — Target 60 Days",
     "58.9% of packaging spend operates on 61-90 day terms, and 5.2% on 121+ days. "
     "ANGLEBOARD (125 days, $329K) and REHRIG (125 days, $105K) have the highest terms. "
     "While extended terms benefit working capital, >90-day terms can strain smaller supplier relationships. "
     "Negotiate a 60-day standard with rebate structure for larger suppliers."),
    ("MEDIUM", "Evaluate Returnable Packaging ROI with Orbis & Rehrig",
     "Orbis ($650K) and Rehrig ($339K) supply plastic top frames. Both companies offer returnable/reusable "
     "container programs. A returnable model typically achieves 15-30% total cost of ownership reduction over "
     "3-5 years. Commission a TCO analysis comparing one-way vs. returnable economics for top frames."),
    ("LOW", "Investigate NCA Maverick Spend",
     "NCA has 25 plant locations but only $1,141,689 in packaging spend — $45,668 per location vs. "
     "EMEA's $223,083 per location (5x gap). This disparity likely indicates off-contract purchasing or "
     "spend flowing through channels outside the AIC system. Survey plant managers and validate against "
     "local purchase orders."),
    ("LOW", "Implement Cat3 Data Governance",
     "18.5% of all 1,010 transactions have no Category 3 classification. For PACKAGING specifically, "
     "20% of rows (6/30) lack sub-categories after correction. Recommend: (1) mandatory Cat3 field at "
     "PO creation, (2) quarterly data completeness audits, (3) a standardized taxonomy guide for "
     "common PACKAGING sub-categories shared with all BUs."),
]
for priority, title, body in recs:
    p = doc.add_paragraph()
    p.add_run(f"[{priority}]  {title}").bold = True
    doc.add_paragraph(body)
    doc.add_paragraph("")

# ── Section 7: Market context
doc.add_heading("7. Market & Supply Chain Context", 1)
for title, body in [
    ("Global Supply Chain Disruptions",
     "Layer pads, wood pallets, and top frames face 8-16+ week lead times. Logistics bottlenecks persist. "
     "Safety stock strategy and supplier lead time SLAs recommended."),
    ("Lumber Price Volatility",
     "Wood Pallet spend ($920,994, 14.6%) is directly exposed to softwood lumber markets. "
     "ISPM-15 heat treatment requirements add cost for international shipments (relevant to Egyptian Pallets)."),
    ("Resin & Polymer Inflation",
     "Stretch wrap (polypropylene) and plastic layer sheets are oil-price sensitive. "
     "Group O's $867K stretch wrap contract should include price escalation caps and review clauses."),
    ("South America Inflation Exposure",
     "SA generates $1.44M (22.9%) in packaging spend from only 3 transactions — highest spend per transaction. "
     "With persistent regional inflation, real costs may be rising faster than reported figures."),
    ("EU Packaging Regulation (PPWR)",
     "EMEA accounts for 52% of spend. Upcoming EU Packaging & Packaging Waste Regulation mandates will "
     "require minimum recycled content and reuse targets, potentially restructuring stretch wrap and "
     "plastic component costs significantly."),
    ("Returnable Packaging Trend",
     "The industry is moving toward circular packaging models. Orbis and Rehrig already support returnable "
     "programs — early adoption could yield competitive TCO advantage and support ESG reporting goals."),
]:
    p = doc.add_paragraph()
    p.add_run(f"{title}: ").bold = True
    p.add_run(body)
doc.add_paragraph("")

# ── Section 8: Stakeholder questions
doc.add_heading("8. Key Stakeholder Interview Questions", 1)
t8 = doc.add_table(rows=1, cols=3)
t8.style = "Table Grid"
bold_row(t8.rows[0])
for i, h in enumerate(["#", "Stakeholder", "Question"]):
    t8.rows[0].cells[i].text = h
for idx, (stake, q) in enumerate([
    ("Plant Operations / Supply Chain", "What are current lead times for pallets, top frames, and layer pads vs. 12 months ago? Have any plants experienced production stoppages due to packaging material shortages?"),
    ("Supply Chain / Planning", "What is the current safety stock policy for packaging materials? Are inventory buffers being held given long lead times?"),
    ("Procurement / Category Lead", "When were the last competitive bids for our top 5 packaging suppliers? Are there active contracts or is this primarily spot buying?"),
    ("Finance / IT", "Can we confirm the EXTENDATA SOLUTIONS transactions ($423K coded to PACKAGING LABELS) — are these IT costs allocated to the wrong category?"),
    ("Finance / Procurement", "What price increases have we absorbed in packaging over the past 2 years? Do supplier agreements include cost escalation or CPI mechanisms?"),
    ("Operations / Logistics", "Are we using returnable packaging programs with Orbis and Rehrig today? What is the current return rate and cost per cycle?"),
    ("NCA Plant Managers", "How is packaging dunnage procured at the plant level — through corporate procurement, a catalog, or local purchasing? Is there spend not captured in AIC?"),
    ("Sustainability / Legal", "Has the company assessed EU PPWR compliance requirements? What changes to our packaging material mix would be required by 2030?"),
], 1):
    r = t8.add_row()
    r.cells[0].text = str(idx)
    r.cells[1].text = stake
    r.cells[2].text = q
doc.add_paragraph("")
doc.add_paragraph("Version 2.0  —  MSBX 5470 Group 2  |  April 2026").runs[0].italic = True

out_docx = r"C:\Users\rbenn\procurement-case-1\MSBX5470_Group2_Recommendations_v2.0.docx"
doc.save(out_docx)
print("DOCX saved:", out_docx)


# ─── POWERPOINT ──────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

DARK_BLUE  = RGBColor(0x1F, 0x38, 0x96)   # deep navy
MID_BLUE   = RGBColor(0x2E, 0x74, 0xB5)
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
ORANGE     = RGBColor(0xED, 0x7D, 0x31)
RED        = RGBColor(0xC0, 0x00, 0x00)
GREEN      = RGBColor(0x37, 0x86, 0x10)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x59, 0x59, 0x59)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank

def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)

def txb(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def bg(slide, color, left=0, top=0, width=13.33, height=7.5):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def rect(slide, color, left, top, width, height):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=14, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, wrap=True):
    tf = txb(slide, left, top, width, height)
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tf

def slide_header(slide, title, subtitle=None):
    # Dark blue top bar
    rect(slide, DARK_BLUE, 0, 0, 13.33, 1.1)
    add_text(slide, title, 0.3, 0.15, 9, 0.75,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.82, 9, 0.35,
                 font_size=13, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
    # Version tag
    add_text(slide, "v2.0", 12.5, 0.2, 0.7, 0.35,
             font_size=12, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

def slide_footer(slide, page_num, total=11):
    rect(slide, DARK_BLUE, 0, 7.15, 13.33, 0.35)
    add_text(slide, "MSBX 5470  |  Group 2  |  Packaging & Dunnage Spend Analysis  |  April 2026",
             0.3, 7.17, 10, 0.28, font_size=9, color=LIGHT_BLUE)
    add_text(slide, f"{page_num} / {total}", 12.3, 7.17, 0.8, 0.28,
             font_size=9, color=WHITE, align=PP_ALIGN.RIGHT)

# ── Slide 1: Title
s = add_slide()
bg(s, DARK_BLUE)
# Accent bar
rect(s, ORANGE, 0, 3.0, 0.12, 2.5)
add_text(s, "PACKAGING & DUNNAGE", 0.5, 1.5, 12, 0.8,
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Spend Analysis & Category Review", 0.5, 2.3, 12, 0.6,
         font_size=26, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "Version 2.0  —  Corrected Dataset  |  3 Confirmed + 5 Suspected Reclassifications", 0.5, 3.0, 12, 0.5,
         font_size=15, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, "MSBX 5470 Procurement & Contracting  |  Group 2  |  April 2026", 0.5, 5.5, 12, 0.5,
         font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "Specialty Aluminum Can Manufacturing", 0.5, 6.0, 12, 0.4,
         font_size=12, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ── Slide 2: What Changed (v2.0 vs v1.0)
s = add_slide()
bg(s, WHITE)
slide_header(s, "What Changed: v1.0 vs. v2.0", "Data quality corrections applied to the PACKAGING category")
slide_footer(s, 2)

# Left box — v1.0
rect(s, LIGHT_GRAY, 0.3, 1.25, 5.8, 5.5)
add_text(s, "Version 1.0  (Original)", 0.5, 1.35, 5.4, 0.45,
         font_size=16, bold=True, color=GRAY)
for i, line in enumerate([
    "Transactions:  33",
    "Total Spend:   $6,373,081.88",
    "Unique Suppliers:  25",
    "",
    "Issues:",
    " SODEXO workwear coded as PACKAGING",
    " Forklift gas coded as PACKAGING",
    " Equipment shaft coded as PACKAGING",
    " 5 additional suspects not flagged",
    " 6 blank Cat3 rows unaddressed",
]):
    add_text(s, line, 0.5, 1.85 + i * 0.38, 5.4, 0.38,
             font_size=13, color=GRAY if not line.startswith(" ") else RED,
             bold=line.startswith("Issues"))

# Right box — v2.0
rect(s, LIGHT_BLUE, 6.5, 1.25, 6.5, 5.5)
add_text(s, "Version 2.0  (Corrected)", 6.7, 1.35, 6.1, 0.45,
         font_size=16, bold=True, color=DARK_BLUE)
for i, line in enumerate([
    "Transactions:  30  (-3 confirmed)",
    "Total Spend:   $6,299,708.55",
    "Unique Suppliers:  22",
    "",
    "Improvements:",
    " 3 mis-classifications removed",
    " 5 more suspects flagged for review",
    " 6 blank Cat3 rows researched",
    " EXTENDATA ($423K) flagged as IT",
    " Payment terms anomalies documented",
]):
    color = DARK_BLUE if not line.startswith(" ") else GREEN
    add_text(s, line, 6.7, 1.85 + i * 0.38, 6.1, 0.38,
             font_size=13, color=color,
             bold=line.startswith("Improvements"))

# Arrow
add_text(s, ">>>", 6.05, 3.7, 0.45, 0.45, font_size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ── Slide 3: Data Quality — Confirmed Reclassifications
s = add_slide()
bg(s, WHITE)
slide_header(s, "Data Quality: 3 Confirmed Reclassifications", "Removed from PACKAGING — corrections applied in workbook (Column L)")
slide_footer(s, 3)

rows = [
    ("Row 543", "SODEXO", "2021 NOV - WORKWEAR", "HR / Uniforms", "$59,402.56", "Largest single error — employee workwear is not packaging"),
    ("Row 755", "JEBLA SA DE CV", "GAS MONTACARGAS 1-15 OCT", "Energy & Utilities", "$10,832.31", "Montacargas = forklift (Spanish); gas = fuel, not packaging"),
    ("Row 211", "ARDAGH METAL BEVERAGE", "SHAFT, 265018", "Equipment - General", "$3,138.46", "Machine shaft is an equipment component, not dunnage"),
]
col_w = [1.1, 2.2, 2.8, 2.0, 1.2, 3.7]
col_x = [0.3, 1.45, 3.7, 6.55, 8.6, 9.85]
headers = ["Row", "Supplier", "PO Description", "Corrected To", "Spend", "Why"]

# Header row
rect(s, DARK_BLUE, 0.3, 1.25, 12.7, 0.45)
for hdr, x, w in zip(headers, col_x, col_w):
    add_text(s, hdr, x, 1.27, w, 0.4, font_size=11, bold=True, color=WHITE)

for i, (row_id, supplier, desc, corrected, spend, why) in enumerate(rows):
    y = 1.75 + i * 1.3
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    rect(s, bg_color, 0.3, y, 12.7, 1.25)
    for val, x, w in zip([row_id, supplier, desc, corrected, spend, why], col_x, col_w):
        color = RED if val == spend else DARK_BLUE
        add_text(s, val, x, y + 0.05, w, 1.15, font_size=11, color=color, wrap=True)

rect(s, ORANGE, 0.3, 5.75, 12.7, 0.6)
add_text(s, "Total removed from PACKAGING:  $73,373.33  |  Corrected spend: $6,299,708.55  (was $6,373,081.88)",
         0.5, 5.82, 12.3, 0.45, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Slide 4: New Suspects
s = add_slide()
bg(s, WHITE)
slide_header(s, "NEW: 5 Additional Suspected Misclassifications", "~$429K still coded to PACKAGING — requires team review and confirmation")
slide_footer(s, 4)

suspects = [
    ("653", "EXTENDATA SOLUTIONS", "SSO FOR BUSINESS USERS", "PACKAGING LABELS", "IT / PROF. SERVICES", "$340,938", "Single Sign-On = IT software. Not a label."),
    ("643", "EXTENDATA SOLUTIONS", "UIM SERVICE RESPONSE",  "PACKAGING LABELS", "IT / PROF. SERVICES", "$57,835",  "UIM service response = IT help desk. Not packaging."),
    ("579", "FOKSAN DOO",          "FOTOKOPIR PAPIR A4",    "PAPER SHEET",       "OFFICE SUPPLIES",     "$16,651",  "Fotokopir = photocopy (Serbian). Office paper."),
    ("551", "NPB AUTOMATION",      "CIRCUIT BREAKER 5A",    "WOODEN FRAME",      "EQUIPMENT",           "$9,938",   "Electrical component coded as wooden frame."),
    ("84",  "GRUPMICROS",          "RENO MANT. MC92XX 1Y",  "PACKAGING LABELS",  "IT / PROF. SERVICES", "$3,710",   "MC92XX = Motorola scanner. Maintenance contract."),
]
col_w2 = [0.7, 2.3, 2.4, 1.8, 1.7, 1.0, 3.0]
col_x2 = [0.3, 1.05, 3.4, 5.85, 7.7, 9.45, 10.5]
headers2 = ["Row", "Supplier", "PO Description", "Current Cat3", "Should Be", "Spend", "Why Suspicious"]

rect(s, DARK_BLUE, 0.3, 1.25, 12.7, 0.45)
for hdr, x, w in zip(headers2, col_x2, col_w2):
    add_text(s, hdr, x, 1.27, w, 0.4, font_size=10, bold=True, color=WHITE)

for i, row_data in enumerate(suspects):
    y = 1.75 + i * 0.98
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    rect(s, bg_color, 0.3, y, 12.7, 0.93)
    for val, x, w in zip(row_data, col_x2, col_w2):
        color = RED if val.startswith("$") else DARK_BLUE
        add_text(s, val, x, y + 0.04, w, 0.85, font_size=10, color=color, wrap=True)

rect(s, RED, 0.3, 6.7, 12.7, 0.55)
add_text(s,
         "Combined potential impact: ~$429,071  |  If confirmed: corrected packaging spend drops to ~$5,870,637",
         0.5, 6.75, 12.3, 0.45, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Slide 5: Executive Summary (corrected numbers)
s = add_slide()
bg(s, WHITE)
slide_header(s, "Executive Summary — Corrected Dataset", "30 transactions, $6.30M spend after 3 confirmed reclassifications")
slide_footer(s, 5)

kpis = [
    ("$6,299,709", "Total Packaging Spend"),
    ("30",         "Transactions"),
    ("22",         "Unique Suppliers"),
    ("$209,990",   "Avg Transaction"),
    ("4.8%",       "of Total Dataset"),
    ("~70 days",   "Avg Payment Terms"),
]
for i, (val, label) in enumerate(kpis):
    col = i % 3
    row_i = i // 3
    x = 0.3 + col * 4.3
    y = 1.3 + row_i * 2.0
    rect(s, DARK_BLUE if row_i == 0 else MID_BLUE, x, y, 4.0, 1.7)
    add_text(s, val, x, y + 0.2, 4.0, 1.0, font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, label, x, y + 1.1, 4.0, 0.5, font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Bottom callout
rect(s, ORANGE, 0.3, 5.4, 12.7, 1.3)
bullets = [
    "Actual spend is 3.2x higher than the original $2M estimate — this category warrants strategic attention",
    "Top 5 suppliers = 73.7% of corrected spend  |  Concentration risk is HIGH",
    "EMEA dominates at 52.1%  |  South America highest spend per transaction ($481,509 avg)",
]
for i, b in enumerate(bullets):
    add_text(s, f"•  {b}", 0.5, 5.48 + i * 0.4, 12.3, 0.38, font_size=12, color=WHITE)

# ── Slide 6: Sub-Category Breakdown
s = add_slide()
bg(s, WHITE)
slide_header(s, "Spend Overview: Sub-Category Breakdown", "Layer Pads is the dominant sub-category at 30.3% of corrected spend")
slide_footer(s, 6)

subcats = [
    ("LAYER PADS",         "$1,909,249", "30.3%", 3.84),
    ("END BAGS",           "$1,071,522", "17.0%", 2.16),
    ("WOOD PALLET",        "$920,994",   "14.6%", 1.85),
    ("STRETCH WRAP",       "$866,966",   "13.8%", 1.74),
    ("TOP FRAMES",         "$539,107",    "8.6%", 1.08),
    ("PACKAGING LABELS*",  "$430,913",    "6.8%", 0.87),
    ("WOODEN FRAME",       "$408,159",    "6.5%", 0.82),
    ("OTHER / BLANK",      "$152,799",    "2.4%", 0.31),
]
bar_colors = [DARK_BLUE, MID_BLUE, DARK_BLUE, MID_BLUE, DARK_BLUE, ORANGE, MID_BLUE, LIGHT_BLUE]

for i, ((label, spend, pct, bar_len), color) in enumerate(zip(subcats, bar_colors)):
    y = 1.25 + i * 0.69
    add_text(s, label, 0.3, y, 2.8, 0.6, font_size=11, bold=True, color=DARK_BLUE)
    if bar_len > 0:
        rect(s, color, 3.15, y + 0.1, bar_len, 0.42)
    add_text(s, spend, 7.1, y, 1.5, 0.6, font_size=11, color=DARK_BLUE, align=PP_ALIGN.RIGHT)
    add_text(s, pct,   8.7, y, 0.8, 0.6, font_size=11, bold=True, color=DARK_BLUE)

add_text(s, "* PACKAGING LABELS includes EXTENDATA SOLUTIONS ($423K) — under review for potential IT reclassification",
         0.3, 6.85, 12.7, 0.35, font_size=10, italic=True, color=RED)

# ── Slide 7: Geographic Analysis
s = add_slide()
bg(s, WHITE)
slide_header(s, "Spend Overview: Geographic Analysis", "EMEA leads at 52.1% — South America has highest spend per transaction")
slide_footer(s, 7)

regions = [
    ("EMEA",      "13", "$3,283,703", "52.1%", "$252,592/txn", 6.62, DARK_BLUE),
    ("SA",        "3",  "$1,444,528", "22.9%", "$481,509/txn", 2.91, MID_BLUE),
    ("NCA",       "9",  "$1,141,689", "18.1%", "$126,854/txn", 2.30, DARK_BLUE),
    ("Corporate", "2",  "$398,772",   "6.3%",  "$199,386/txn", 0.80, MID_BLUE),
    ("Aerosol",   "3",  "$31,016",    "0.5%",  "$10,339/txn",  0.06, LIGHT_BLUE),
]
for i, (bu, txns, spend, pct, per_txn, bar_len, color) in enumerate(regions):
    y = 1.3 + i * 1.1
    rect(s, color, 0.3, y, 0.6, 0.9)
    add_text(s, bu, 1.0, y + 0.1, 1.8, 0.7, font_size=13, bold=True, color=DARK_BLUE)
    add_text(s, spend, 2.85, y + 0.1, 1.8, 0.7, font_size=13, color=DARK_BLUE, align=PP_ALIGN.RIGHT)
    if bar_len > 0.05:
        rect(s, color, 4.75, y + 0.2, bar_len, 0.5)
    add_text(s, pct, 4.85 + bar_len, y + 0.1, 1.0, 0.7, font_size=13, bold=True, color=DARK_BLUE)
    add_text(s, per_txn, 9.8, y + 0.1, 2.0, 0.7, font_size=11, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)

# Callout box
rect(s, LIGHT_GRAY, 0.3, 6.5, 12.7, 0.7)
add_text(s,
         "NCA Flag:  25 plant locations, only $45,668/location vs. EMEA $252,592/location — "
         "potential maverick spend or data gaps in NCA region",
         0.5, 6.55, 12.3, 0.6, font_size=12, italic=True, color=RED)

# ── Slide 8: Supplier Analysis (Pareto)
s = add_slide()
bg(s, WHITE)
slide_header(s, "Supplier Analysis: Pareto & Concentration Risk", "Top 5 suppliers = 73.7% of spend — single-source risk on Layer Pads & Stretch Wrap")
slide_footer(s, 8)

top_supps = [
    (1,  "ANGLEBOARD",              "$1,459,249", "23.2%", "39.2%",  "Layer Pads",         4.67),
    (2,  "GSD VERPACKUNGEN",        "$1,011,306", "16.1%", "52.9%",  "End Bags",           3.24),
    (3,  "GROUP O",                 "$866,966",   "13.8%", "63.4%",  "Stretch Wrap",       2.77),
    (4,  "EGYPTIAN PALLETS",        "$657,912",   "10.4%", "73.7%",  "Wood Pallets",       2.11),
    (5,  "ORBIS",                   "$650,000",   "10.3%", "80.5%",  "Top Frames",         2.08),
    (6,  "EXTENDATA SOLUTIONS *",   "$423,277",    "6.7%", "86.8%",  "Labels (review)",    1.35),
    (7,  "FINE ARTS DE MEXICO",     "$398,221",    "6.3%", "92.2%",  "Wooden Frame",       1.28),
    (8,  "REHRIG",                  "$339,107",    "5.4%", "94.0%",  "Top Frames",         1.09),
]
ys = 1.25
rect(s, DARK_BLUE, 0.3, ys, 12.7, 0.42)
for x, w, h in [(0.35, 0.5, "#"), (0.9, 2.8, "Supplier"), (3.75, 1.3, "Spend"),
                 (5.1, 0.8, "% of Pkg"), (5.95, 0.9, "Cumul."), (6.9, 1.7, "Product"), (8.65, 4.3, "Bar")]:
    add_text(s, h, x, ys+0.03, w, 0.35, font_size=10, bold=True, color=WHITE)

for i, (rank, supp, spend, pct, cumul, product, bar_len) in enumerate(top_supps):
    y = 1.7 + i * 0.62
    bg_c = LIGHT_GRAY if i % 2 == 0 else WHITE
    rect(s, bg_c, 0.3, y, 12.7, 0.6)
    bar_color = ORANGE if rank <= 5 else MID_BLUE
    if i < 5:
        bar_color = DARK_BLUE
    if rank == 6:
        bar_color = RED
    for val, x, w in [
        (str(rank), 0.35, 0.5),
        (supp,      0.9,  2.75),
        (spend,     3.75, 1.25),
        (pct,       5.1,  0.78),
        (cumul,     5.95, 0.85),
        (product,   6.9,  1.65),
    ]:
        color = RED if rank == 6 and val not in (str(rank), spend, pct, cumul) else DARK_BLUE
        if val == spend:
            color = MID_BLUE
        add_text(s, val, x, y+0.08, w, 0.44, font_size=10, color=color)
    if bar_len > 0:
        rect(s, bar_color, 8.65, y+0.12, bar_len, 0.36)

add_text(s, "* EXTENDATA SOLUTIONS flagged as potential IT spend — classification under review",
         0.3, 6.8, 12.7, 0.35, font_size=10, italic=True, color=RED)

# ── Slide 9: Recommendations
s = add_slide()
bg(s, WHITE)
slide_header(s, "Procurement Recommendations", "Prioritized action plan based on corrected spend analysis")
slide_footer(s, 9)

recs_slide = [
    ("CRITICAL", RED,       "Audit EXTENDATA SOLUTIONS ($423K)", "SSO and UIM service responses coded to PACKAGING LABELS. Escalate to Finance/IT. Could reduce packaging budget by 6.7%."),
    ("HIGH",     ORANGE,    "Confirm 4 Additional Misclassifications", "Circuit breaker, copy paper, scanner maintenance contract = ~$30,300 still in PACKAGING. Review PO documentation."),
    ("HIGH",     ORANGE,    "Assign Cat3 to 6 Blank Rows ($97,463)", "Kartonfabrik surcharge ($47K), ULINE products, MULTI SERVICES items need sub-category assignment."),
    ("HIGH",     ORANGE,    "Address Supplier Concentration Risk", "Top 5 = 73.7% of spend. Angleboard is sole-source for layer pads ($1.46M). Run competitive bids; develop secondary sources."),
    ("MEDIUM",   MID_BLUE,  "Consolidate Wood Pallet Supply (5 Suppliers)", "5 vendors, $921K. Consolidate to 2 preferred. Retain Egyptian Pallets as primary; qualify one secondary."),
    ("MEDIUM",   MID_BLUE,  "Standardize Payment Terms to 60 Days", "Angleboard & Rehrig at 125 days ($329K). Target 60-day standard with rebate for early payment."),
    ("MEDIUM",   MID_BLUE,  "Evaluate Returnable Packaging ROI", "Orbis + Rehrig ($989K) offer returnable programs. Commission TCO analysis — 15-30% potential savings."),
    ("LOW",      GRAY,      "Investigate NCA Maverick Spend", "$45,668/location vs EMEA $252,592. Survey plant managers for off-contract purchasing."),
]
for i, (priority, color, title, body) in enumerate(recs_slide):
    y = 1.22 + i * 0.74
    rect(s, color, 0.3, y, 1.1, 0.65)
    add_text(s, priority, 0.31, y+0.1, 1.08, 0.45, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, 1.5, y+0.01, 4.5, 0.35, font_size=11, bold=True, color=DARK_BLUE)
    add_text(s, body, 1.5, y+0.35, 11.3, 0.35, font_size=10, color=GRAY)

# ── Slide 10: Market Context
s = add_slide()
bg(s, WHITE)
slide_header(s, "Market Research & Supply Chain Context", "Key risk factors shaping the Packaging & Dunnage category")
slide_footer(s, 10)

market = [
    ("Supply Chain Disruptions", "HIGH",   ORANGE, "Wood pallets, top frames, layer pads — 8-16+ week lead times persist. Safety stock strategy required."),
    ("Lumber Price Volatility",   "HIGH",   ORANGE, "Wood Pallet spend ($921K) directly exposed. ISPM-15 heat treatment adds cross-border costs."),
    ("Resin / Polymer Inflation", "HIGH",   ORANGE, "Stretch wrap & plastic pallets tied to oil prices. Group O's $867K contract needs escalation clauses."),
    ("SA Inflation Exposure",     "HIGH",   ORANGE, "SA = $1.44M (22.9%), only 3 txns — highest risk per dollar. Regional inflation erodes real contract value."),
    ("EU PPWR Regulations",       "MEDIUM", MID_BLUE, "EMEA (52% of spend) faces EU packaging mandates on recycled content and reuse targets by 2030."),
    ("Returnable Packaging",      "OPPORTUNITY", GREEN, "Orbis & Rehrig offer returnable models. TCO savings of 15-30% achievable — supports ESG targets."),
    ("Supplier Concentration",    "HIGH",   ORANGE, "Single-source risk on layer pads (Angleboard) and stretch wrap (Group O) — production continuity at risk."),
]
for i, (topic, impact, color, detail) in enumerate(market):
    col = i % 2
    row_i = i // 2
    x = 0.3 + col * 6.5
    y = 1.3 + row_i * 1.5
    if i == 6:
        x = 3.55
        y = 1.3 + 3 * 1.5
    rect(s, color, x, y, 6.0, 1.35)
    add_text(s, topic, x+0.12, y+0.08, 4.0, 0.45, font_size=13, bold=True, color=WHITE)
    add_text(s, f"[{impact}]", x+4.2, y+0.08, 1.6, 0.45, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(s, detail, x+0.12, y+0.55, 5.7, 0.75, font_size=10, color=WHITE)

# ── Slide 11: Stakeholder Questions + Next Steps
s = add_slide()
bg(s, WHITE)
slide_header(s, "Stakeholder Questions & Next Steps", "Interview guide and immediate action items")
slide_footer(s, 11)

questions = [
    ("Ops / Supply Chain", "What are current lead times for top packaging materials vs. 12 months ago?"),
    ("Procurement",        "When were last competitive bids run for top 5 packaging suppliers?"),
    ("Finance / IT",       "Confirm EXTENDATA SOLUTIONS ($423K) — packaging or IT costs?"),
    ("Finance",            "What price increases absorbed in packaging last 2 years? Escalation clauses in contracts?"),
    ("Operations",         "Returnable packaging with Orbis & Rehrig — current return rate and TCO?"),
    ("NCA Plant Mgrs",     "How is packaging procured at the plant level — centrally or locally?"),
    ("Legal / Sustain.",   "EU PPWR compliance assessment — what packaging changes are required by 2030?"),
]
add_text(s, "Stakeholder Interview Questions", 0.3, 1.2, 7.8, 0.45, font_size=14, bold=True, color=DARK_BLUE)
for i, (stake, q) in enumerate(questions):
    y = 1.7 + i * 0.7
    rect(s, LIGHT_BLUE, 0.3, y, 1.5, 0.6)
    add_text(s, stake, 0.32, y+0.08, 1.46, 0.45, font_size=9, bold=True, color=DARK_BLUE)
    add_text(s, q, 1.9, y+0.08, 6.1, 0.5, font_size=11, color=DARK_BLUE)

# Next steps
rect(s, DARK_BLUE, 8.5, 1.2, 4.6, 5.95)
add_text(s, "Next Steps", 8.65, 1.3, 4.3, 0.45, font_size=14, bold=True, color=WHITE)
next_steps = [
    ("1", "Escalate EXTENDATA flag to Finance/IT"),
    ("2", "Confirm 4 additional mis-class. rows"),
    ("3", "Assign Cat3 to 6 blank rows"),
    ("4", "Run competitive bid — Layer Pads"),
    ("5", "Run competitive bid — Stretch Wrap"),
    ("6", "Request TCO analysis from Orbis/Rehrig"),
    ("7", "Survey NCA plant managers re: local buying"),
    ("8", "Validate SA supplier contracts vs. inflation"),
]
for i, (num, step) in enumerate(next_steps):
    rect(s, ORANGE, 8.65, 1.85 + i * 0.64, 0.4, 0.55)
    add_text(s, num, 8.66, 1.9 + i * 0.64, 0.38, 0.44, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, step, 9.12, 1.9 + i * 0.64, 3.8, 0.5, font_size=11, color=WHITE)

out_pptx = r"C:\Users\rbenn\procurement-case-1\MSBX5470_Group2_CaseStudy_v2.0.pptx"
prs.save(out_pptx)
print("PPTX saved:", out_pptx)
print("Done.")
