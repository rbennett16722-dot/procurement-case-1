#!/usr/bin/env python3
"""
generate_presentation.py
Builds the PowerPoint deck for MSBX 5470 Case 1: Packaging-Dunnage Category Review
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Color Palette ─────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1F, 0x4E, 0x79)
MED_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
ORANGE     = RGBColor(0xED, 0x7D, 0x31)
RED        = RGBColor(0xC0, 0x00, 0x00)
AMBER      = RGBColor(0xFF, 0xC0, 0x00)
GREEN      = RGBColor(0x70, 0xAD, 0x47)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x40, 0x40, 0x40)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY   = RGBColor(0xA0, 0xA0, 0xA0)
HIGHLIGHT  = RGBColor(0xDD, 0xEB, 0xF7)

# ── Presentation Setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helper Functions ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=12, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.0, fill=DARK_BLUE)
    add_text(slide, title, 0.3, 0.06, 12.5, 0.58, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.60, 12.5, 0.38, size=12, color=LIGHT_BLUE)
    add_rect(slide, 0, 7.35, 13.33, 0.15, fill=ORANGE)


def set_cell(cell, text, size=10, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.CENTER, fill_color=None):
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    cell.text = str(text)
    for para in cell.text_frame.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def make_table(slide, data, l, t, w, h, col_widths=None, font_size=10):
    rows = len(data)
    cols = len(data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(l), Inches(t),
                                 Inches(w), Inches(h)).table
    if col_widths:
        for c, cw in enumerate(col_widths):
            tbl.columns[c].width = Inches(cw)
    return tbl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
add_rect(s1, 0, 5.65, 13.33, 0.14, fill=ORANGE)
add_rect(s1, 0, 7.15, 13.33, 0.35, fill=MED_BLUE)

add_text(s1, "PACKAGING-DUNNAGE", 1.0, 1.1, 11.33, 1.2,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "CATEGORY REVIEW", 1.0, 2.1, 11.33, 1.1,
         size=44, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s1, "Specialty Aluminum Can Manufacturing", 1.0, 3.25, 11.33, 0.6,
         size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s1, "Prepared for: Chief Procurement Officer & Chief Financial Officer",
         1.0, 3.75, 11.33, 0.5, size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s1, "January 2023  |  2022 Spend Analysis & Market Research",
         1.0, 4.2, 11.33, 0.5, size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s1,
         "Presented by:  [Team Member 1]   \u2022   [Team Member 2]   \u2022   [Team Member 3]   \u2022   [Team Member 4]",
         1.0, 4.95, 11.33, 0.5, size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_header(s2, "Executive Summary",
           "Actual spend is 3.2x the $2M estimate \u2014 this category demands immediate strategic attention")

# 6 KPI boxes (2 rows x 3 cols)
kpis = [
    ("$6.37M",  "Total 2022 Category Spend",      DARK_BLUE),
    ("+218%",   "Above $2M Initial CPO Estimate",  RED),
    ("25",      "Unique Suppliers (33 Transactions)", MED_BLUE),
    ("72.9%",   "Spend in Top 5 Suppliers",         RED),
    ("52.5%",   "Spend Concentrated in EMEA",       DARK_BLUE),
    ("4.85%",   "Share of Total Company Spend",     MED_BLUE),
]
kpi_positions = [
    (0.3,  1.1), (4.5,  1.1), (8.7,  1.1),
    (0.3,  2.4), (4.5,  2.4), (8.7,  2.4),
]
for (val, lbl, clr), (lp, tp) in zip(kpis, kpi_positions):
    add_rect(s2, lp, tp, 3.9, 1.15, fill=HIGHLIGHT, line_color=clr, line_width=2)
    add_text(s2, val, lp+0.1, tp+0.05, 3.7, 0.65,
             size=26, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_text(s2, lbl, lp+0.1, tp+0.68, 3.7, 0.44,
             size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Key findings header bar
add_rect(s2, 0.3, 3.65, 12.7, 0.36, fill=DARK_BLUE)
add_text(s2, "KEY FINDINGS", 0.45, 3.68, 5.0, 0.3,
         size=11, bold=True, color=WHITE)

bullets = [
    "Layer Pads + End Bags + Wood Pallets = 60% of spend \u2014 the 3 sub-categories most exposed to active supply chain disruption and raw material price volatility",
    "SINGLE-SOURCE RISK: Angleboard ($1.46M, Layer Pads) and Group O ($867K, Stretch Wrap) have no identified backup supplier \u2014 any disruption could halt production",
    "NCA anomaly: 25 US locations at only $45.7K/location vs. $180.6K in South America \u2014 investigate possible maverick spend or untracked purchasing",
]
for i, b in enumerate(bullets):
    tp = 4.07 + i * 0.66
    add_rect(s2, 0.3, tp, 0.08, 0.55, fill=ORANGE)
    add_text(s2, b, 0.5, tp+0.03, 12.5, 0.58, size=10.5, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Category Background
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_header(s3, "Category Background",
           "Packaging-Dunnage: What We Buy, Where We Use It, and Why It Matters Now")

# Left: sub-category table
add_rect(s3, 0.3, 1.05, 6.1, 0.42, fill=MED_BLUE)
add_text(s3, "WHAT WE BUY  (10 Sub-Categories)", 0.4, 1.08, 5.9, 0.37,
         size=11, bold=True, color=WHITE)

subcats = [
    ("Layer Pads",             "$1,909,249",  "30.0%"),
    ("End Bags / Paper Bags",  "$1,071,522",  "16.8%"),
    ("Wood Pallets",           "$920,994",    "14.5%"),
    ("Stretch Wrap / Film",    "$866,966",    "13.6%"),
    ("Top Frames",             "$539,107",    "8.5%"),
    ("Packaging Labels",       "$430,913",    "6.8%"),
    ("Wooden Frames",          "$408,159",    "6.4%"),
    ("Strap Tape + Other",     "$226,172",    "3.5%"),
]
for i, (name, spend, pct) in enumerate(subcats):
    tp = 1.55 + i * 0.56
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(s3, 0.3, tp, 6.1, 0.54, fill=bg)
    add_text(s3, name, 0.42, tp+0.07, 3.5, 0.42, size=11, color=DARK_GRAY)
    add_text(s3, spend, 3.85, tp+0.07, 1.3, 0.42, size=11, bold=True,
             color=DARK_BLUE, align=PP_ALIGN.RIGHT)
    add_text(s3, pct, 5.2, tp+0.07, 1.1, 0.42, size=11, color=MED_BLUE,
             align=PP_ALIGN.RIGHT)

add_rect(s3, 0.3, 6.15, 6.1, 0.54, fill=DARK_BLUE)
add_text(s3, "TOTAL CATEGORY SPEND", 0.42, 6.2, 3.5, 0.42,
         size=11, bold=True, color=WHITE)
add_text(s3, "$6,373,082", 3.85, 6.2, 1.3, 0.42, size=11, bold=True,
         color=ORANGE, align=PP_ALIGN.RIGHT)
add_text(s3, "100%", 5.2, 6.2, 1.1, 0.42, size=11, bold=True,
         color=WHITE, align=PP_ALIGN.RIGHT)

# Right: context
add_rect(s3, 6.7, 1.05, 6.3, 0.42, fill=MED_BLUE)
add_text(s3, "COMPANY CONTEXT & SITUATION", 6.8, 1.08, 6.1, 0.37,
         size=11, bold=True, color=WHITE)

contexts = [
    ("Global Footprint",
     "48 Locations: 25 NCA  |  15 EMEA  |  8 South America",
     DARK_BLUE),
    ("Supply Disruptions",
     "8-16+ week lead times for pallets, top frames, and layer pads. Port congestion and logistics bottlenecks persist post-COVID.",
     RED),
    ("Inflation Pressure",
     "SA inflation >90% (Argentina). EMEA energy costs rising. Lumber and resin prices volatile globally.",
     AMBER),
    ("Workforce Turnover",
     "2 US locations experiencing high turnover \u2014 creates demand variability and potential maverick buying risk.",
     ORANGE),
    ("Key Named Suppliers",
     "Orbis and Rehrig identified as larger/strategic suppliers. Both offer returnable packaging programs.",
     MED_BLUE),
    ("Initial CPO Estimate",
     "Supervisor estimated $2M category spend. Actual: $6.37M \u2014 a $4.37M / 218% variance requiring immediate attention.",
     RED),
]
for i, (title, detail, accent) in enumerate(contexts):
    tp = 1.57 + i * 0.92
    add_rect(s3, 6.7, tp, 0.09, 0.82, fill=accent)
    add_text(s3, title, 6.87, tp+0.02, 6.0, 0.34, size=11, bold=True, color=DARK_GRAY)
    add_text(s3, detail, 6.87, tp+0.36, 6.0, 0.52, size=9.5, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Spend by Sub-Category (Bar Chart)
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_header(s4, "Spend Overview: Sub-Category Breakdown",
           "Top 4 sub-categories = 74.9% of spend | All 4 face active supply disruption or raw material inflation")

cd = ChartData()
cd.categories = ['Layer Pads', 'End Bags', 'Wood Pallet',
                 'Stretch Wrap', 'Top Frames', 'Pkg Labels',
                 'Wooden Frame', 'Other']
cd.add_series('Spend ($)', (
    1909249, 1071522, 920994,
    866966, 539107, 430913,
    408159, 226172
))

chart4 = s4.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.3), Inches(1.1), Inches(7.8), Inches(6.1),
    cd
).chart
chart4.has_legend = False
chart4.has_title = False
chart4.plots[0].gap_width = 70
chart4.series[0].format.fill.solid()
chart4.series[0].format.fill.fore_color.rgb = MED_BLUE
chart4.value_axis.number_format = '$#,##0'

# Callout boxes
callouts4 = [
    ("74.9%",  "of spend in top 4\nsub-categories",            RED),
    ("$2.98M", "Layer Pads + End Bags\n(supply disruption risk)", DARK_BLUE),
    ("$1.84M", "Wood Pallets + Wooden Frames\n(lumber exposure)", ORANGE),
    ("$867K",  "Stretch Wrap \u2014 single-source\nsupplier (Group O)", AMBER),
]
for i, (val, lbl, clr) in enumerate(callouts4):
    tp = 1.15 + i * 1.52
    add_rect(s4, 8.25, tp, 4.8, 1.4, fill=LIGHT_GRAY, line_color=clr, line_width=2)
    add_text(s4, val, 8.35, tp+0.04, 4.6, 0.72, size=28, bold=True, color=clr,
             align=PP_ALIGN.CENTER)
    add_text(s4, lbl, 8.35, tp+0.74, 4.6, 0.62, size=11, color=DARK_GRAY,
             align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Spend by Geography
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_header(s5, "Spend Overview: Geographic Analysis",
           "EMEA dominates at 52.5% | South America shows highest spend per location | NCA spend gap flagged")

# Pie chart - geography
geo_cd = ChartData()
geo_cd.categories = ['EMEA', 'South America', 'NCA', 'Corporate', 'Aerosol']
geo_cd.add_series('Spend', (3346244, 1444528, 1141689, 398772, 41848))

chart5 = s5.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    Inches(0.3), Inches(1.1), Inches(6.2), Inches(5.5),
    geo_cd
).chart
from pptx.enum.chart import XL_LEGEND_POSITION
chart5.has_legend = True
chart5.legend.position = XL_LEGEND_POSITION.BOTTOM
chart5.has_title = False

# Regional table - right side
add_rect(s5, 6.7, 1.1, 6.4, 0.42, fill=DARK_BLUE)
geo_rows = [
    ["Region",         "Spend",       "% Total", "Loca-\ntions", "$/Location"],
    ["EMEA",           "$3,346,244",  "52.5%",   "15",     "$223,083"],
    ["South America",  "$1,444,528",  "22.7%",   "8",      "$180,566"],
    ["NCA",            "$1,141,689",  "17.9%",   "25",     "$45,668  ⚠"],
    ["Corporate",      "$398,772",    "6.3%",    "—",      "—"],
    ["Aerosol",        "$41,848",     "0.7%",    "—",      "—"],
    ["TOTAL",          "$6,373,082",  "100%",    "48",     "—"],
]
geo_cwids = [1.7, 1.5, 1.1, 0.9, 1.2]
geo_tbl = make_table(s5, geo_rows, 6.7, 1.1, 6.4, 3.35,
                     col_widths=geo_cwids)
for r, row_data in enumerate(geo_rows):
    for c, val in enumerate(row_data):
        cell = geo_tbl.cell(r, c)
        is_header = (r == 0)
        is_total  = (r == len(geo_rows) - 1)
        is_nca    = (r == 3)
        fc = WHITE if (is_header or is_total) else DARK_GRAY
        bg = DARK_BLUE if is_header else (MED_BLUE if is_total else
             (AMBER if is_nca else (LIGHT_GRAY if r % 2 == 0 else WHITE)))
        set_cell(cell, val, size=10,
                 bold=(is_header or is_total),
                 color=fc, fill_color=bg,
                 align=PP_ALIGN.CENTER)

# Notes
notes5 = [
    ("\u26a0  NCA ANOMALY",
     "25 US locations at $45.7K/location is 75% below SA per-location spend. "
     "Possible maverick spend or transactions coded outside PACKAGING category.",
     AMBER),
    ("\u26a0  SA INFLATION RISK",
     "Only 3 transactions but $1.44M spend. Argentina inflation >90% in 2022. "
     "Real cost impact likely understated \u2014 watch 2023 closely.",
     RED),
    ("\u2714  EMEA STRATEGIC FOCUS",
     "$3.35M (52.5%) of spend. One location experiencing energy disruption. "
     "Largest region by spend warrants priority supplier relationship management.",
     MED_BLUE),
]
for i, (hdr, body, clr) in enumerate(notes5):
    tp = 4.6 + i * 0.92
    add_rect(s5, 6.7, tp, 0.09, 0.8, fill=clr)
    add_text(s5, hdr, 6.87, tp+0.02, 6.1, 0.3, size=10, bold=True, color=clr)
    add_text(s5, body, 6.87, tp+0.32, 6.1, 0.55, size=9.5, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Supplier Pareto & Concentration Risk
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_header(s6, "Supplier Analysis: Pareto & Concentration Risk",
           "6 of 25 suppliers = 80% of spend | Top 2 are SINGLE-SOURCE | This is a risk, not a strength")

# Supplier table (left)
sup_rows = [
    ["#", "Supplier",              "Spend",       "% Total", "Cum. %", "Product"],
    ["1", "Angleboard",            "$1,459,249",  "22.9%",   "22.9%",  "Layer Pads \u26a0"],
    ["2", "GSD Verpackungen",      "$1,011,306",  "15.9%",   "38.8%",  "End Bags"],
    ["3", "Group O",               "$866,966",    "13.6%",   "52.4%",  "Stretch Wrap \u26a0"],
    ["4", "Egyptian Pallets",      "$657,912",    "10.3%",   "62.7%",  "Wood Pallets"],
    ["5", "Orbis",                 "$650,000",    "10.2%",   "72.9%",  "Plastic Pallets/Frames"],
    ["6", "Extendata Solutions",   "$423,277",    "6.6%",    "79.5%",  "Pkg Labels"],
    ["7", "Fine Arts de Mexico",   "$398,221",    "6.2%",    "85.8%",  "Wooden Frames"],
    ["8", "Rehrig",                "$339,107",    "5.3%",    "91.1%",  "Top Frames"],
    ["9", "Aldhana Wood Ind.",     "$118,791",    "1.9%",    "93.0%",  "Wood Pallets"],
    ["10","Jay Wood Industry",     "$79,535",     "1.2%",    "94.2%",  "Wood/Top Frames"],
    ["\u2014","Other 15 Suppliers","$368,718",   "5.8%",    "100%",   "Various"],
]
sup_cwids = [0.38, 2.05, 1.3, 0.95, 0.95, 1.72]
sup_tbl = make_table(s6, sup_rows, 0.3, 1.05, 7.35, 6.35,
                     col_widths=sup_cwids)
pareto_rows = {1, 2, 3, 4, 5, 6}  # rows that together reach 80%
for r, row_data in enumerate(sup_rows):
    for c, val in enumerate(row_data):
        cell = sup_tbl.cell(r, c)
        is_hdr   = (r == 0)
        is_last  = (r == len(sup_rows) - 1)
        is_80    = (r in pareto_rows)
        is_ss    = (r in {1, 3})  # single source
        bg = (DARK_BLUE if is_hdr else
              (LIGHT_GRAY if is_last else
               (RGBColor(0xFF, 0xE0, 0xE0) if is_ss else
                (HIGHLIGHT if is_80 else WHITE))))
        fc = WHITE if is_hdr else DARK_GRAY
        align = PP_ALIGN.LEFT if c in {1, 5} else PP_ALIGN.CENTER
        set_cell(cell, val, size=9, bold=is_hdr, color=fc,
                 fill_color=bg, align=align)

# Single-source label
add_rect(s6, 0.3, 1.57, 0.09, 0.52, fill=RED)
add_rect(s6, 0.3, 2.44, 0.09, 0.52, fill=RED)

# Column chart (right)
par_cd = ChartData()
par_cd.categories = ['Angleboard', 'GSD Verpack.', 'Group O',
                     'Egyptian Pal.', 'Orbis', 'Extendata',
                     'Fine Arts MX', 'Rehrig', 'Aldhana', 'Jay Wood', 'Other 15']
par_cd.add_series('Spend ($)', (
    1459249, 1011306, 866966, 657912,
    650000, 423277, 398221, 339107,
    118791, 79535, 368718
))
chart6 = s6.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(7.75), Inches(1.1), Inches(5.3), Inches(4.7),
    par_cd
).chart
chart6.has_legend = False
chart6.has_title = False
chart6.value_axis.number_format = '$#,##0'
chart6.plots[0].gap_width = 50
chart6.series[0].format.fill.solid()
chart6.series[0].format.fill.fore_color.rgb = MED_BLUE

# Bottom callouts
callouts6 = [
    ("6 of 25 Suppliers = 80% Spend",   ORANGE),
    ("Single-Source Risk = $2.33M",     RED),
    ("Orbis + Rehrig Returnable Opp.",  GREEN),
    ("Top 10 of 25 = 94.2% Spend",     DARK_BLUE),
]
for i, (txt, clr) in enumerate(callouts6):
    lp = 7.75 + (i % 2) * 2.7
    tp = 6.0  + (i // 2) * 0.62
    add_rect(s6, lp, tp, 2.55, 0.55, fill=LIGHT_GRAY, line_color=clr, line_width=2)
    add_text(s6, txt, lp+0.1, tp+0.07, 2.35, 0.42, size=10, bold=True, color=clr,
             align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Market Research Findings
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_header(s7, "Market Research Findings",
           "10 trends identified | 4 HIGH impact | 4 MEDIUM impact | 2 OPPORTUNITIES")

mkt_rows = [
    ["#", "Issue / Market Trend", "Impact", "Affected Sub-Categories"],
    ["1",  "Global Supply Chain Disruptions \u2014 8-16+ week lead times for pallets, top frames, layer pads. Port congestion and logistics bottlenecks persist post-COVID.", "HIGH", "Pallets, Top Frames, Layer Pads"],
    ["2",  "Resin & Raw Material Price Volatility \u2014 Stretch wrap and plastic pallets tied to polypropylene/PE resin and oil prices. Direct cost pass-through from suppliers.", "HIGH", "Stretch Wrap, Plastic Pallets, Layer Pads"],
    ["3",  "Lumber Price Volatility \u2014 Wood pallet costs ($921K, 14.5%) directly exposed to softwood lumber markets. ISPM-15 heat treatment adds cost for international shipments.", "HIGH", "Wood Pallets, Wooden Frames"],
    ["4",  "South American Inflation Exposure \u2014 $1.44M in SA with Argentina inflation >90%. Real cost increases likely understated in 2022 reported spend.", "HIGH", "All SA-Sourced Materials"],
    ["5",  "European Energy Crisis Impact \u2014 Energy disruption at 1 EMEA location. Energy-intensive manufacturers (plastic, corrugated) passing cost increases to buyers.", "MEDIUM", "Layer Pads, End Bags, Stretch Wrap"],
    ["6",  "Sustainability & Regulatory Pressure \u2014 EU Packaging & Packaging Waste Regulation (PPWR) mandates recycled content. Potential material substitutions and compliance costs ahead.", "MEDIUM", "Stretch Wrap, Plastic Pallets"],
    ["7",  "Workforce Turnover at 2 US Sites \u2014 High turnover creates demand volatility, packaging waste, and maverick buying risk. Procurement process gaps possible.", "MEDIUM", "All Sub-Categories"],
    ["8",  "Data Visibility & Classification Gaps \u2014 9 transactions ($171K) unclassified in source data. NCA spend low for 25 locations. Category visibility is incomplete.", "MEDIUM", "All Sub-Categories"],
    ["9",  "Returnable Packaging Opportunity \u2014 Orbis + Rehrig (combined $989K) both offer returnable/reusable pallet programs. Industry TCO savings: 15-30% vs. one-way pallets.", "OPPTY", "Pallets, Top Frames"],
    ["10", "Global Volume Aggregation \u2014 Regional purchasing silos likely paying different prices to the same global suppliers. Centralizing could unlock significant volume leverage.", "OPPTY", "All Suppliers"],
]
impact_fill = {"HIGH": RED, "MEDIUM": AMBER, "OPPTY": GREEN}
mkt_cwids = [0.35, 7.5, 1.1, 3.75]
mkt_tbl = make_table(s7, mkt_rows, 0.3, 1.04, 12.7, 6.3,
                     col_widths=mkt_cwids)
for r, row_data in enumerate(mkt_rows):
    for c, val in enumerate(row_data):
        cell = mkt_tbl.cell(r, c)
        is_hdr = (r == 0)
        impact = row_data[2] if r > 0 else None
        bg = (DARK_BLUE if is_hdr else
              (impact_fill.get(impact, WHITE) if c == 2 else
               (LIGHT_GRAY if r % 2 == 0 else WHITE)))
        fc = WHITE if (is_hdr or c == 2) else DARK_GRAY
        align = PP_ALIGN.CENTER if c in {0, 2} else PP_ALIGN.LEFT
        set_cell(cell, val, size=(10 if is_hdr else 9),
                 bold=(is_hdr or c == 2), color=fc,
                 fill_color=bg, align=align)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Risk Analysis & Mitigation
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_header(s8, "Risk Analysis & Mitigation Strategy",
           "6 active risks | 4 rated HIGH | Supplier concentration is the most immediate and quantifiable risk")

risk_rows = [
    ["Risk",                                    "Category",           "Impact", "Like-\nlihood", "Mitigation Strategy"],
    ["Single-Source Supplier Exposure",          "Supply Continuity",  "HIGH",   "MEDIUM",
     "Dual-source Angleboard (Layer Pads, $1.46M) and Group O (Stretch Wrap, $867K) within Q1-Q2 2023. "
     "Issue mini-RFQ to 2-3 qualified alternates per category. Combined at-risk spend: $2.33M."],
    ["South American Inflation / FX Volatility", "Financial",          "HIGH",   "HIGH",
     "Implement index-linked pricing clauses tied to regional inflation indices. "
     "Review SA supplier contracts quarterly. Explore USD-denominated contracts to reduce FX exposure on $1.44M SA spend."],
    ["Lumber & Resin Price Volatility",          "Financial",          "HIGH",   "HIGH",
     "Add commodity escalation/de-escalation clauses tied to PPI Lumber Index (wood pallets, $921K) "
     "and ICIS resin pricing (stretch wrap, $867K). Hedge commodity exposure where volume warrants."],
    ["EMEA Energy Disruption",                   "Supply Continuity",  "HIGH",   "MEDIUM",
     "Assess energy surcharge pass-through clauses in EU supplier contracts. "
     "Identify backup supply sources outside energy-disrupted regions for top EMEA suppliers ($3.35M spend)."],
    ["Spend Data Gaps / Maverick Spend",          "Compliance",         "MEDIUM", "HIGH",
     "Audit NCA plant purchasing immediately. Mandate Category 3 on all POs. "
     "Investigate NCA spend gap ($45.7K vs. $180.6K/location in SA). "
     "Implement spend visibility dashboard with monthly refresh."],
    ["EU PPWR Regulatory Compliance",             "Regulatory",         "MEDIUM", "HIGH",
     "Engage EMEA suppliers on recycled content roadmap by Q3 2023. "
     "Pilot sustainable alternatives at 1-2 EMEA locations. "
     "Assess full PPWR compliance cost in category budget planning for 2024."],
]
risk_cwids = [2.2, 1.75, 0.9, 0.95, 6.9]
risk_tbl = make_table(s8, risk_rows, 0.3, 1.04, 12.7, 6.3,
                      col_widths=risk_cwids)
for r, row_data in enumerate(risk_rows):
    for c, val in enumerate(row_data):
        cell = risk_tbl.cell(r, c)
        is_hdr = (r == 0)
        imp  = row_data[2] if r > 0 else None
        lkly = row_data[3] if r > 0 else None
        bg = (DARK_BLUE if is_hdr else
              (impact_fill.get(imp,    WHITE) if c == 2 else
               (impact_fill.get(lkly, WHITE) if c == 3 else
                (LIGHT_GRAY if r % 2 == 0 else WHITE))))
        fc = WHITE if (is_hdr or c in {2, 3}) else DARK_GRAY
        align = PP_ALIGN.CENTER if c in {2, 3} else PP_ALIGN.LEFT
        set_cell(cell, val, size=(10 if is_hdr else 9),
                 bold=(is_hdr or c in {2, 3}), color=fc,
                 fill_color=bg, align=align)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Stakeholder Interview Questions
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_header(s9, "Stakeholder Interview Questions",
           "12 questions across 3 functional groups \u2014 grounded in spend data findings")

groups9 = [
    {
        "title":  "OPERATIONS & SUPPLY CHAIN",
        "color":  MED_BLUE,
        "qs": [
            ("Q1",  "What are current lead times for pallets, top frames, and layer pads? How have they changed vs. 12 months ago?",
             "Plant Ops / Supply Chain"),
            ("Q2",  "Have any plants had production stoppages due to packaging-dunnage shortages in the last 12 months?",
             "Plant Managers"),
            ("Q3",  "What is the current inventory policy (days on hand)? Are we safety-stocking given 8-16 week lead times?",
             "Supply Chain / Planning"),
            ("Q4",  "Are we using returnable/reusable packaging programs with Orbis and Rehrig? What is the return rate and TCO?",
             "Operations / Logistics"),
        ]
    },
    {
        "title":  "PROCUREMENT & COMMERCIAL",
        "color":  ORANGE,
        "qs": [
            ("Q5",  "When were the last competitive bids for our top 5 suppliers? Do we have contracts or is this spot buying?",
             "Procurement / Category Lead"),
            ("Q6",  "Why is 72.9% of spend with 5 suppliers? Is concentration by strategic design or lack of sourcing events?",
             "Global Procurement Director"),
            ("Q7",  "Is there maverick spend outside procurement? NCA at $45.7K/location seems very low for 25 US sites.",
             "Plant Managers / Finance"),
            ("Q8",  "Is there opportunity to aggregate global volume across regions with Orbis / Rehrig for better pricing?",
             "Global Procurement Director"),
        ]
    },
    {
        "title":  "FINANCE & STRATEGY",
        "color":  GREEN,
        "qs": [
            ("Q9",  "What price increases have we absorbed over 2 years? Do our supplier agreements include escalation clauses?",
             "Finance / Procurement"),
            ("Q10", "Are packaging costs passed to customers or absorbed internally? What is the COGS and margin impact?",
             "CFO / Commercial"),
            ("Q11", "Has the company evaluated EU PPWR compliance requirements? What operational changes are needed in EMEA?",
             "Legal / Sustainability / EMEA"),
            ("Q12", "Who approves packaging specs, and how often are they reviewed for cost optimization opportunities?",
             "Engineering / Quality"),
        ]
    },
]

for ci, grp in enumerate(groups9):
    lp = 0.3 + ci * 4.35
    add_rect(s9, lp, 1.05, 4.1, 0.42, fill=grp["color"])
    add_text(s9, grp["title"], lp+0.05, 1.07, 4.0, 0.38,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for qi, (qnum, question, stakeholder) in enumerate(grp["qs"]):
        tp = 1.55 + qi * 1.47
        bg = LIGHT_GRAY if qi % 2 == 0 else WHITE
        add_rect(s9, lp, tp, 4.1, 1.43, fill=bg)
        add_rect(s9, lp, tp, 0.52, 1.43, fill=grp["color"])
        add_text(s9, qnum, lp+0.04, tp+0.5, 0.46, 0.45,
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s9, question, lp+0.58, tp+0.05, 3.45, 1.0, size=9.5, color=DARK_GRAY)
        add_text(s9, f"\u2192 {stakeholder}", lp+0.58, tp+1.13, 3.45, 0.28,
                 size=8, italic=True, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Recommendations & Next Steps
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
add_header(s10, "Recommendations & Next Steps",
           "5 specific actions | Est. $583K\u2013$1.37M in value opportunity | Prioritized by risk and impact")

recs10 = [
    ("01", "IMMEDIATE", RED,
     "Dual-Source Layer Pads & Stretch Wrap",
     "$1.46M (Angleboard) and $867K (Group O) are single-source with no identified backup. "
     "Issue mini-RFQ to 2-3 qualified alternates per category by end of Q1 2023. "
     "Combined exposure protected: $2.33M.",
     "Risk Reduction\n$2.33M Protected"),
    ("02", "SHORT-TERM", ORANGE,
     "Run Competitive Sourcing Events for Top 5 Suppliers",
     "Unclear when RFPs were last conducted. A 5% cost reduction on $4.71M (top 5) = $235K savings. "
     "A 10% reduction = $471K. Prioritize Angleboard and GSD Verpackungen as first targets.",
     "Savings Est.\n$235K\u2013$471K/yr"),
    ("03", "SHORT-TERM", ORANGE,
     "Evaluate Returnable Packaging TCO with Orbis and Rehrig",
     "Orbis + Rehrig combined spend = $989K (15.5% of category). Both offer returnable pallet programs. "
     "Industry TCO savings: 15\u201330% vs. one-way pallets. Commission a formal business case in Q2 2023.",
     "TCO Savings Est.\n$148K\u2013$297K/yr"),
    ("04", "MEDIUM-TERM", MED_BLUE,
     "Add Commodity Escalation Clauses & Audit NCA Spend",
     "Lumber (14.5%) and resin (13.6%) expose the company to volatile commodity markets with no contractual protection. "
     "Simultaneously, investigate why NCA has only $45.7K/location vs. $180.6K in SA \u2014 "
     "potential untracked spend or PO coding errors.",
     "Cost Control\n+ Data Integrity"),
    ("05", "MEDIUM-TERM", MED_BLUE,
     "Launch Global Volume Aggregation Strategy",
     "Suppliers like Orbis and Rehrig likely operate across all regions. "
     "Centralizing $6.37M in category spend under a global master agreement could unlock "
     "volume-based pricing discounts and standardized terms across all 48 locations.",
     "Strategic Value\n$300K\u2013$600K/yr est."),
]

for i, (num, priority, clr, title, detail, value) in enumerate(recs10):
    tp = 1.08 + i * 1.27
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(s10, 0.3, tp, 12.7, 1.23, fill=bg)
    add_rect(s10, 0.3, tp, 0.72, 1.23, fill=clr)
    add_text(s10, num, 0.3, tp+0.35, 0.72, 0.55,
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s10, 1.1, tp, 1.55, 0.3, fill=clr)
    add_text(s10, priority, 1.12, tp+0.04, 1.51, 0.26,
             size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s10, title, 1.1, tp+0.3, 8.9, 0.4, size=13, bold=True, color=DARK_BLUE)
    add_text(s10, detail, 1.1, tp+0.68, 8.9, 0.52, size=9.5, color=DARK_GRAY)
    add_rect(s10, 10.15, tp+0.12, 2.8, 1.0, fill=clr)
    add_text(s10, value, 10.2, tp+0.22, 2.7, 0.8,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "Case1_Packaging_Dunnage_Category_Review.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
for i, sl in enumerate(prs.slides, 1):
    print(f"  Slide {i:2d}: OK")
